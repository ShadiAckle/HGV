#!/usr/bin/env python3
"""Generate HGV Compensation Hub — Data Model & MCP stakeholder deck."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

OUT = Path(__file__).resolve().parents[1] / "docs" / "HGV_Comp_Data_Model_and_MCP.pptx"

# HGV / luxury hospitality palette
NAVY = RGBColor(0x0A, 0x25, 0x40)
NAVY_MID = RGBColor(0x12, 0x36, 0x5A)
TEAL = RGBColor(0x14, 0xB8, 0xA6)
GOLD = RGBColor(0xC9, 0xA2, 0x27)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
SLATE = RGBColor(0x64, 0x74, 0x8B)
LIGHT = RGBColor(0xF1, 0xF5, 0xF9)
CHARCOAL = RGBColor(0x1E, 0x29, 0x3B)


def set_bg(slide, color: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_bar(slide, top=Inches(0), height=Inches(0.12), color=TEAL) -> None:
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), top, Inches(10), height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def add_footer(slide, text: str = "HGV Compensation Hub  |  Confidential") -> None:
    box = slide.shapes.add_textbox(Inches(0.5), Inches(7.05), Inches(9), Inches(0.35))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(9)
    p.font.color.rgb = SLATE
    p.alignment = PP_ALIGN.CENTER


def add_title_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, NAVY)
    add_bar(slide, Inches(0), Inches(0.14), TEAL)

    t = slide.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(8.5), Inches(1.2))
    tf = t.text_frame
    p = tf.paragraphs[0]
    p.text = "HGV Compensation Hub"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = WHITE

    p2 = tf.add_paragraph()
    p2.text = "Data Model & Agent Integration (MCP)"
    p2.font.size = Pt(26)
    p2.font.color.rgb = TEAL
    p2.space_before = Pt(12)

    sub = slide.shapes.add_textbox(Inches(0.75), Inches(3.6), Inches(8), Inches(1))
    stf = sub.text_frame
    sp = stf.paragraphs[0]
    sp.text = "Governed semantic layer over your Cognos & PwC ETL  •  AI-ready MCP tools"
    sp.font.size = Pt(16)
    sp.font.color.rgb = RGBColor(0xCB, 0xD5, 0xE1)

    meta = slide.shapes.add_textbox(Inches(0.75), Inches(6.2), Inches(5), Inches(0.5))
    mp = meta.text_frame.paragraphs[0]
    mp.text = "edw_dev_hris.hgv_comp  |  Databricks Apps  |  Model Context Protocol"
    mp.font.size = Pt(11)
    mp.font.color.rgb = SLATE


def add_section(prs: Presentation, title: str, subtitle: str = "") -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, NAVY_MID)
    add_bar(slide, Inches(3.2), Inches(0.08), TEAL)

    t = slide.shapes.add_textbox(Inches(0.75), Inches(2.5), Inches(8.5), Inches(1))
    tf = t.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = WHITE

    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(18)
        p2.font.color.rgb = TEAL
        p2.space_before = Pt(16)


def add_content_slide(
    prs: Presentation,
    title: str,
    bullets: list[str],
    dark: bool = False,
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, NAVY if dark else WHITE)

    title_color = WHITE if dark else NAVY
    body_color = RGBColor(0xE2, 0xE8, 0xF0) if dark else CHARCOAL

    t = slide.shapes.add_textbox(Inches(0.6), Inches(0.45), Inches(8.8), Inches(0.8))
    tp = t.text_frame.paragraphs[0]
    tp.text = title
    tp.font.size = Pt(28)
    tp.font.bold = True
    tp.font.color.rgb = title_color

    accent = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.6), Inches(1.15), Inches(1.2), Inches(0.06))
    accent.fill.solid()
    accent.fill.fore_color.rgb = TEAL
    accent.line.fill.background()

    body = slide.shapes.add_textbox(Inches(0.75), Inches(1.45), Inches(8.7), Inches(5.4))
    tf = body.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.size = Pt(15 if not dark else 14)
        p.font.color.rgb = body_color
        p.space_after = Pt(10)
        p.line_spacing = 1.15

    add_footer(slide, "HGV Compensation Hub  |  Confidential" if not dark else "HGV Compensation Hub  |  Confidential — Dark")


def add_two_column_slide(prs: Presentation, title: str, left_title: str, left_items: list[str], right_title: str, right_items: list[str]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, WHITE)

    t = slide.shapes.add_textbox(Inches(0.6), Inches(0.45), Inches(8.8), Inches(0.7))
    t.text_frame.paragraphs[0].text = title
    t.text_frame.paragraphs[0].font.size = Pt(28)
    t.text_frame.paragraphs[0].font.bold = True
    t.text_frame.paragraphs[0].font.color.rgb = NAVY

    for col, (ctitle, items, left) in enumerate([(left_title, left_items, 0.6), (right_title, right_items, 5.1)]):
        card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(left), Inches(1.35), Inches(4.2), Inches(5.5))
        card.fill.solid()
        card.fill.fore_color.rgb = LIGHT
        card.line.color.rgb = RGBColor(0xE2, 0xE8, 0xF0)

        hdr = slide.shapes.add_textbox(Inches(left + 0.2), Inches(1.55), Inches(3.8), Inches(0.5))
        hp = hdr.text_frame.paragraphs[0]
        hp.text = ctitle
        hp.font.size = Pt(16)
        hp.font.bold = True
        hp.font.color.rgb = NAVY_MID

        body = slide.shapes.add_textbox(Inches(left + 0.2), Inches(2.05), Inches(3.8), Inches(4.5))
        tf = body.text_frame
        tf.word_wrap = True
        for i, item in enumerate(items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = f"• {item}"
            p.font.size = Pt(12)
            p.font.color.rgb = CHARCOAL
            p.space_after = Pt(6)

    add_footer(slide)


def add_table_slide(prs: Presentation, title: str, headers: list[str], rows: list[list[str]], col_widths: list[float] | None = None) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, WHITE)

    t = slide.shapes.add_textbox(Inches(0.6), Inches(0.45), Inches(8.8), Inches(0.7))
    t.text_frame.paragraphs[0].text = title
    t.text_frame.paragraphs[0].font.size = Pt(26)
    t.text_frame.paragraphs[0].font.bold = True
    t.text_frame.paragraphs[0].font.color.rgb = NAVY

    cols = len(headers)
    rows_n = len(rows) + 1
    table = slide.shapes.add_table(rows_n, cols, Inches(0.5), Inches(1.25), Inches(9), Inches(0.4 * rows_n)).table

    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = Inches(w)

    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        p = cell.text_frame.paragraphs[0]
        p.font.bold = True
        p.font.size = Pt(11)
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    for r, row in enumerate(rows, 1):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.text = val
            if r % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(10)
            p.font.color.rgb = CHARCOAL
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    add_footer(slide)


def add_architecture_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, WHITE)

    t = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(8.8), Inches(0.7))
    t.text_frame.paragraphs[0].text = "Architecture: Your Data → Governed Layer → Experiences"
    t.text_frame.paragraphs[0].font.size = Pt(26)
    t.text_frame.paragraphs[0].font.bold = True
    t.text_frame.paragraphs[0].font.color.rgb = NAVY

    layers = [
        ("HGV Source Systems (existing ETL)", NAVY, [
            "edw_dev_cognos.cognos_fm — it_smt_detail, marketing, personnel",
            "edw_dev_hris.pwcmodels — commissions (PwC / Varicent)",
        ]),
        ("Semantic Layer  edw_dev_hris.hgv_comp", NAVY_MID, [
            "Views + materialized Delta tables (FY2025–FY2026 window)",
            "No direct agent access to billion-row Cognos tables",
        ]),
        ("Consumption", TEAL, [
            "Compensation Hub web app (persona dashboards)",
            "MCP tools + ResponsesAgent (AI Playground, Supervisor, bots)",
        ]),
    ]

    y = 1.35
    for label, color, lines in layers:
        box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(1.2), Inches(y), Inches(7.6), Inches(1.35))
        box.fill.solid()
        box.fill.fore_color.rgb = color
        box.line.fill.background()

        tb = slide.shapes.add_textbox(Inches(1.45), Inches(y + 0.12), Inches(7.1), Inches(1.1))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = WHITE
        for line in lines:
            lp = tf.add_paragraph()
            lp.text = line
            lp.font.size = Pt(11)
            lp.font.color.rgb = RGBColor(0xE2, 0xE8, 0xF0)
            lp.space_before = Pt(4)

        if y < 5.5:
            arrow = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.DOWN_ARROW, Inches(4.65), Inches(y + 1.38), Inches(0.5), Inches(0.35))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = GOLD
            arrow.line.fill.background()
        y += 1.85

    add_footer(slide)


def add_mcp_flow_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, NAVY)
    add_bar(slide, Inches(0), Inches(0.1), TEAL)

    t = slide.shapes.add_textbox(Inches(0.6), Inches(0.45), Inches(8.8), Inches(0.7))
    t.text_frame.paragraphs[0].text = "MCP Integration Flow"
    t.text_frame.paragraphs[0].font.size = Pt(28)
    t.text_frame.paragraphs[0].font.bold = True
    t.text_frame.paragraphs[0].font.color.rgb = WHITE

    steps = [
        ("1", "Orchestrator", "AI Playground, Agent Bricks Supervisor,\nLangGraph, enterprise copilot"),
        ("2", "POST /mcp", "Streamable HTTP — tool discovery\n& JSON-schema validated calls"),
        ("3", "Comp Hub Server", "Same SQL + Claude path as the web UI\n(no duplicate logic)"),
        ("4", "Unity Catalog", "edw_dev_hris.hgv_comp\n(governed reads only)"),
        ("5", "Response", "Structured JSON back to the agent\nfor narration or workflow"),
    ]

    x = 0.45
    for num, head, desc in steps:
        card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(1.6), Inches(1.75), Inches(4.8))
        card.fill.solid()
        card.fill.fore_color.rgb = NAVY_MID
        card.line.color.rgb = TEAL

        nb = slide.shapes.add_textbox(Inches(x + 0.15), Inches(1.8), Inches(0.5), Inches(0.4))
        nb.text_frame.paragraphs[0].text = num
        nb.text_frame.paragraphs[0].font.size = Pt(22)
        nb.text_frame.paragraphs[0].font.bold = True
        nb.text_frame.paragraphs[0].font.color.rgb = TEAL

        hb = slide.shapes.add_textbox(Inches(x + 0.15), Inches(2.35), Inches(1.45), Inches(0.8))
        hp = hb.text_frame.paragraphs[0]
        hp.text = head
        hp.font.size = Pt(13)
        hp.font.bold = True
        hp.font.color.rgb = WHITE

        db = slide.shapes.add_textbox(Inches(x + 0.15), Inches(3.2), Inches(1.45), Inches(2.8))
        dp = db.text_frame.paragraphs[0]
        dp.text = desc
        dp.font.size = Pt(10)
        dp.font.color.rgb = RGBColor(0xCB, 0xD5, 0xE1)
        db.text_frame.word_wrap = True

        x += 1.88

    note = slide.shapes.add_textbox(Inches(0.6), Inches(6.55), Inches(8.8), Inches(0.5))
    note.text_frame.paragraphs[0].text = "Deployed as mcp-hgv-comp-hub  •  Auth: Databricks OAuth or AGENT_API_KEY"
    note.text_frame.paragraphs[0].font.size = Pt(11)
    note.text_frame.paragraphs[0].font.color.rgb = SLATE
    note.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER


def add_closing_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, NAVY)
    add_bar(slide, Inches(3.3), Inches(0.1), TEAL)

    t = slide.shapes.add_textbox(Inches(0.75), Inches(2.4), Inches(8.5), Inches(1))
    tf = t.text_frame
    p = tf.paragraphs[0]
    p.text = "One governed comp layer."
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    p2 = tf.add_paragraph()
    p2.text = "Web UI, AI agents, and MCP — same truth."
    p2.font.size = Pt(22)
    p2.font.color.rgb = TEAL
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(14)

    q = slide.shapes.add_textbox(Inches(0.75), Inches(5.2), Inches(8.5), Inches(0.6))
    qp = q.text_frame.paragraphs[0]
    qp.text = "Questions & next steps"
    qp.font.size = Pt(18)
    qp.font.color.rgb = RGBColor(0xCB, 0xD5, 0xE1)
    qp.alignment = PP_ALIGN.CENTER


def main() -> None:
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    add_title_slide(prs)

    add_content_slide(
        prs,
        "Executive Summary",
        [
            "We built a governed compensation semantic layer on top of HGV’s existing Cognos and PwC ETL — not a parallel data mart.",
            "The HGV Compensation Hub gives marketing reps, managers, finance, and admins persona-specific views of the same underlying truth.",
            "AI is grounded in warehouse SQL context — never hallucinated payouts or tour counts.",
            "Model Context Protocol (MCP) exposes that same governed data to Databricks AI Playground, Supervisor Agents, and external orchestrators.",
            "Production target: edw_dev_hris.hgv_comp on HGV VDI Databricks with FY2025–FY2026 materialized marketing core for interactive speed.",
        ],
        dark=True,
    )

    add_section(prs, "Part 1", "Data Model & Your Sources")

    add_content_slide(
        prs,
        "The Challenge",
        [
            "Cognos tour tables (it_smt_detail) contain billions of rows — direct app queries time out on VDI.",
            "PwC commissions data is wide and granular — unsuitable as an ad-hoc analytics surface for frontline reps.",
            "Demo personas and synthetic seeds don’t reflect live marketing reps or real tour history.",
            "AI copilots without governance risk inventing compensation answers.",
            "We need one semantic layer that is fast, auditable, and reusable across UI and agents.",
        ],
    )

    add_architecture_slide(prs)

    add_two_column_slide(
        prs,
        "What We Read From HGV (Not Replaced)",
        "Cognos — edw_dev_cognos.cognos_fm",
        [
            "it_smt_detail — tours, show/qualify, volume",
            "it_smt_marketing — channel, office, program",
            "it_smt_personnel — rep assignment per tour",
            "it_smt_contract / it_uni_contract — deal linkage",
            "it_uni_lead — guest / lead master",
        ],
        "PwC / Varicent — edw_dev_hris.pwcmodels",
        [
            "commissions — field rep payouts & participants",
            "Pay dates drive dim_period calendar",
            "Business unit, title, plan metadata",
            "60-month lookback in views; FY2025–26 in materialized core",
        ],
    )

    add_table_slide(
        prs,
        "Semantic Layer: edw_dev_hris.hgv_comp",
        ["Layer", "Examples", "Role"],
        [
            ["Dimensions", "dim_rep, dim_marketing_rep, dim_period, dim_team, dim_guest", "Who, when, where — rep pickers & filters"],
            ["Marketing facts", "fact_marketing_tour_payout, fact_marketing_rep_period", "Tour ledger & period KPIs for My Comp"],
            ["Field sales facts", "fact_payout, fact_quota_attainment, fact_deal_credit", "Quota, earnings, deal credits"],
            ["Manager / finance", "fact_team_snapshot, fact_tour_quality, dim_finance_period", "Team rollups & finance corridors"],
            ["Reference / scenario", "industry_comp_benchmark, scenario_run, plan_assessment_*", "Benchmarks & what-if (writable)"],
            ["Staging (internal)", "_src_tour_spine, _src_rep_directory", "Bounded Cognos joins — not exposed to app"],
        ],
        col_widths=[1.4, 3.8, 3.3],
    )

    add_content_slide(
        prs,
        "Production Strategy: Views → Materialization",
        [
            "Script 12/15 — Live views over Cognos/PwC with lookback windows (36mo tours, 60mo commissions).",
            "Correct for logic & governance; too slow for interactive VDI when queried on every page load.",
            "Script 16 — One-time CTAS to Delta tables for marketing core (FY2025–FY2026 only).",
            "Materialized: dim_marketing_rep, fact_marketing_tour_payout, fact_marketing_rep_period, dim_period, dim_rep.",
            "Refresh cadence: re-run script 16 after ETL loads (weekly/monthly). App reads local Delta in seconds.",
        ],
        dark=True,
    )

    add_section(prs, "Part 2", "Model Context Protocol (MCP)")

    add_content_slide(
        prs,
        "What Is MCP — In Plain Language?",
        [
            "MCP (Model Context Protocol) is an open standard for giving AI agents structured tools — like APIs the model can call safely.",
            "Instead of every team wiring custom HTTP to the comp warehouse, orchestrators discover tools once and pass validated JSON arguments.",
            "Our MCP server runs inside the Databricks App (mcp-hgv-comp-hub) — same process as the REST API, no extra infrastructure.",
            "Agents never get raw SQL access — only pre-approved tools that query hgv_comp.",
            "Same Claude-backed comp agent and warehouse paths as the web UI — one implementation, many consumers.",
        ],
    )

    add_mcp_flow_slide(prs)

    add_table_slide(
        prs,
        "Six MCP Tools — What We Surface",
        ["Tool", "What the agent receives", "HGV use case"],
        [
            ["ask_comp_agent", "Grounded natural-language answer (Claude + warehouse context)", "“Why is my Q2 payout lower than Q1?”"],
            ["get_marketing_workspace", "Full My Comp payload: KPIs, tours, money map, chargebacks", "Manager or bot reviewing rep performance"],
            ["get_tour_context", "Guest 360 + comp impact for one tour", "Drill-down on a specific arrival / tour"],
            ["get_comp_metadata", "Reps, teams, periods, scenarios directory", "Populate agent context & pickers"],
            ["search_comp_entities", "Typeahead matches for reps, teams, deals", "@mention resolution in copilots"],
            ["hub_health", "Liveness & version check", "Ops monitoring & Playground registration"],
        ],
        col_widths=[1.8, 3.5, 3.2],
    )

    add_two_column_slide(
        prs,
        "Governance & Security",
        "What we enforce",
        [
            "Unity Catalog permissions on edw_dev_hris + edw_dev_cognos",
            "No arbitrary SQL from MCP — tools only",
            "Lookback windows & FY2025–26 materialization caps scan size",
            "Production mode skips demo persona bootstrap",
            "Audit trail via fact_comp_admin_log (writable Delta)",
        ],
        "How HGV connects",
        [
            "AI Playground → MCP Servers → mcp-hgv-comp-hub",
            "Agent Bricks Supervisor → add MCP tool",
            "External bots → POST /mcp + AGENT_API_KEY",
            "Catalog discovery → GET /api/agent/info",
            "VDI local demo → vdi-start.ps1 on port 8000",
        ],
    )

    add_content_slide(
        prs,
        "What’s Next",
        [
            "Complete script 16 materialization on VDI (FY2025–FY2026) and validate live marketing rep picker.",
            "Expand materialization window or schedule nightly refresh job after ETL.",
            "Register mcp-hgv-comp-hub in AI Gateway for Playground & Supervisor demos.",
            "Stakeholder walkthrough: My Comp, Team Performance, Finance Intelligence personas.",
            "Optional: fine-grained MCP scopes per persona (marketing vs finance tools).",
        ],
        dark=True,
    )

    add_closing_slide(prs)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(f"Wrote {OUT} ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
